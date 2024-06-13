import asyncio
import json
import logging
import os
import argparse
from pptx import Presentation
from openai import AsyncOpenAI, APIError, APIConnectionError, RateLimitError

# Set up logging
logging.basicConfig(level=logging.INFO)

# Load API key from environment variable
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")

if OPENAI_API_KEY is None:
    logging.error("OPENAI_API_KEY is not set. Please set it in your environment variables.")
    exit(1)

# Create OpenAI client
client = AsyncOpenAI(api_key=OPENAI_API_KEY)

# Function to extract text from each slide
def extract_text_from_slide(slide):
    text = ""
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            text += shape.text + " "
    return text.strip()

# Asynchronous function to get explanation from OpenAI with exponential backoff and rate limiting
async def get_explanation(client, text, retries=5):
    delay = 1  # Initial delay in seconds
    for attempt in range(retries):
        try:
            response = await client.chat.completions.create(
                model="gpt-3.5-turbo",
                messages=[
                    {"role": "system", "content": "You are a helpful assistant."},
                    {"role": "user", "content": f"Explain the following presentation slide content succinctly:\n\n{text}"}
                ],
                max_tokens=1500,
                temperature=0.5,
            )
            logging.info(f"Received explanation: {response.choices[0].message['content'].strip()[:60]}...")
            return response.choices[0].message['content'].strip()
        except RateLimitError as e:
            logging.warning("Rate limit exceeded. Waiting for 1 minute before retrying...")
            await asyncio.sleep(60)  # Wait for 1 minute before retrying
        except APIConnectionError as e:
            logging.error(f"Failed to connect to OpenAI API: {str(e)}. Retrying in {delay} seconds...")
            await asyncio.sleep(delay)
            delay = min(delay * 2, 60)  # Exponential backoff
        except APIError as e:
            logging.error(f"OpenAI API returned an error: {str(e)}. Retrying in {delay} seconds...")
            await asyncio.sleep(delay)
            delay = min(delay * 2, 60)  # Exponential backoff
        except Exception as e:
            logging.error(f"Error processing slide: {str(e)}")
            return f"Error processing slide: {str(e)}"
    return "Failed to get explanation after several retries due to rate limits."

# Main asynchronous function to process the presentation
async def process_presentation(file_path):
    logging.info(f"Processing presentation: {file_path}")
    ppt = Presentation(file_path)
    all_texts = [extract_text_from_slide(slide) for slide in ppt.slides if extract_text_from_slide(slide)]

    explanations = []
    for i, text in enumerate(all_texts):
        if text:
            logging.info(f"Submitting slide {i+1}/{len(all_texts)} for explanation...")
            explanation = await get_explanation(client, text)
            explanations.append(explanation)
            if (i + 1) % 3 == 0:
                logging.info("Processed 3 slides. Waiting for 1 minute to comply with rate limits...")
                await asyncio.sleep(60)  # Wait for 1 minute after processing 3 slides
            else:
                await asyncio.sleep(20)  # Ensure at least 20 seconds between requests to avoid exceeding rate limits
    return explanations

# Save explanations to a JSON file
def save_explanations(file_path, explanations):
    output_file = f"{os.path.splitext(file_path)[0]}.json"
    logging.info(f"Saving explanations to {output_file}")
    with open(output_file, 'w') as f:
        json.dump(explanations, f, indent=4)

# CLI functionality
def main():
    parser = argparse.ArgumentParser(description="Explain PowerPoint slides using GPT-3.5")
    parser.add_argument("file_path", type=str, help="Path to the PowerPoint file")
    args = parser.parse_args()

    explanations = asyncio.run(process_presentation(args.file_path))
    save_explanations(args.file_path, explanations)
    logging.info(f"Explanations saved to {os.path.splitext(args.file_path)[0]}.json")

if __name__ == "__main__":
    main()
